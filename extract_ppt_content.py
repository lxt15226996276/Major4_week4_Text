from pptx import Presentation
import os

ppt_dir = r"d:\Project\UnityProject\Major4_week4_Text\Tests\UnityKnowledgeMap\ppt"

def extract_ppt(ppt_path):
    prs = Presentation(ppt_path)
    content = []
    for i, slide in enumerate(prs.slides):
        slide_info = {"slide_num": i + 1, "title": "", "text_blocks": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs = []
                for p in shape.text_frame.paragraphs:
                    text = p.text.strip()
                    if text:
                        paragraphs.append(text)
                if paragraphs:
                    full_text = "\n".join(paragraphs)
                    if i == 0 or ("单元" in full_text) or ("Unit" in full_text) or ("第十一" in full_text) or ("第十二" in full_text) or ("第十三" in full_text) or ("第十四" in full_text) or ("第十五" in full_text):
                        slide_info["title"] = full_text
                    else:
                        slide_info["text_blocks"].append(full_text)
        content.append(slide_info)
    return content

for filename in os.listdir(ppt_dir):
    if filename.endswith(".pptx"):
        ppt_path = os.path.join(ppt_dir, filename)
        print(f"\n{'='*60}")
        print(f"文件: {filename}")
        print('='*60)
        content = extract_ppt(ppt_path)
        for slide in content:
            print(f"\n--- 幻灯片 {slide['slide_num']} ---")
            if slide['title']:
                print(f"标题: {slide['title']}")
            for block in slide['text_blocks']:
                print(f"内容:\n{block}\n")